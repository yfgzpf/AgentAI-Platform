"""
OpenClaw 核心服务
整合工具管理器、技能管理器、工具执行器的统一服务
"""
import os
import json
import asyncio
import webbrowser
import subprocess
import datetime
from typing import Dict, Any, Optional, List
from pathlib import Path
from loguru import logger


class OpenClawCore:
    """OpenClaw 核心服务 - 统一管理工具和技能"""
    
    def __init__(self, skills_dir: str = None):
        self.skills_dir = skills_dir or os.path.join(os.path.dirname(__file__), '..', '..', '..', 'skills')
        
        self.tools = {
            'local': {
                'ocr': {
                    'name': '本地OCR工具',
                    'description': '使用本地OCR库识别图片中的文字',
                    'capabilities': ['text_recognition'],
                    'available': True
                },
                'image_editing': {
                    'name': '本地图像编辑工具',
                    'description': '使用本地库进行图像处理和编辑',
                    'capabilities': ['crop', 'resize', 'filter', 'rotate'],
                    'available': True
                },
                'file_operations': {
                    'name': '本地文件操作工具',
                    'description': '执行本地文件的创建、读取、更新和删除操作',
                    'capabilities': ['read', 'write', 'copy', 'delete'],
                    'available': True
                },
                'desktop_automation': {
                    'name': '本地桌面自动化工具',
                    'description': '执行本地桌面自动化操作',
                    'capabilities': ['open_app', 'close_app', 'click', 'type', 'screenshot'],
                    'available': True
                },
                'browser_automation': {
                    'name': '浏览器自动化工具',
                    'description': '使用Playwright进行浏览器自动化',
                    'capabilities': ['navigate', 'click', 'type', 'screenshot', 'extract'],
                    'available': True
                }
            },
            'external': {
                'qianwen_vl': {
                    'name': '千问VL模型',
                    'description': '使用千问VL模型进行图像识别和分析',
                    'capabilities': ['text_recognition', 'image_analysis', 'object_detection'],
                    'available': True,
                    'cost_per_use': 0.01
                },
                'deepseek_chat': {
                    'name': 'DeepSeek对话模型',
                    'description': '使用DeepSeek模型进行对话和文本生成',
                    'capabilities': ['chat', 'text_generation', 'summarization', 'translation'],
                    'available': True,
                    'cost_per_use': 0.005
                },
                'image_generation': {
                    'name': 'AI图像生成',
                    'description': '使用AI模型生成图像',
                    'capabilities': ['image_generation', 'style_transfer'],
                    'available': True,
                    'cost_per_use': 0.05
                },
                'video_generation': {
                    'name': 'AI视频生成',
                    'description': '使用豆包Seedance生成视频',
                    'capabilities': ['video_generation'],
                    'available': True,
                    'cost_per_use': 0.1
                }
            }
        }
        
        self.skills: Dict[str, Any] = {}
        self.skill_registry: Dict[str, Dict[str, Any]] = {}
        self.tool_usage_history: List[Dict[str, Any]] = []
        self.execution_cache: Dict[str, Any] = {}
        
        logger.info("OpenClaw 核心服务初始化完成")
    
    def get_available_tools(self, capability: Optional[str] = None) -> Dict[str, Any]:
        """获取可用的工具"""
        available_tools = {'local': [], 'external': []}
        
        for tool_name, tool_info in self.tools['local'].items():
            if tool_info['available']:
                if not capability or capability in tool_info['capabilities']:
                    available_tools['local'].append({'name': tool_name, 'info': tool_info})
        
        for tool_name, tool_info in self.tools['external'].items():
            if tool_info['available']:
                if not capability or capability in tool_info['capabilities']:
                    available_tools['external'].append({'name': tool_name, 'info': tool_info})
        
        return {"success": True, "tools": available_tools}
    
    def select_best_tool(self, capability: str, requirements: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """选择最适合的工具"""
        for tool_name, tool_info in self.tools['local'].items():
            if tool_info['available'] and capability in tool_info['capabilities']:
                return {"success": True, "tool_type": "local", "tool_name": tool_name, "tool_info": tool_info}
        
        best_external_tool = None
        lowest_cost = float('inf')
        
        for tool_name, tool_info in self.tools['external'].items():
            if tool_info['available'] and capability in tool_info['capabilities']:
                cost = tool_info.get('cost_per_use', 0)
                if cost < lowest_cost:
                    lowest_cost = cost
                    best_external_tool = {"tool_name": tool_name, "tool_info": tool_info}
        
        if best_external_tool:
            return {"success": True, "tool_type": "external", **best_external_tool}
        
        return {"success": False, "message": f"没有可用的工具满足能力要求: {capability}"}
    
    def execute_tool(self, tool_id: str, **kwargs) -> Dict[str, Any]:
        """执行工具"""
        try:
            logger.info(f"执行工具: {tool_id}, 参数: {kwargs}")
            
            method_name = f"_execute_{tool_id}"
            if hasattr(self, method_name):
                method = getattr(self, method_name)
                result = method(**kwargs)
                self._record_tool_usage(tool_id, kwargs, result)
                return result
            else:
                return {"success": False, "message": f"未找到工具 {tool_id} 的执行方法"}
        
        except Exception as e:
            logger.error(f"执行工具 {tool_id} 失败: {str(e)}")
            return {"success": False, "message": f"执行失败: {str(e)}"}
    
    def _execute_open_browser(self, url: str = "https://www.baidu.com") -> Dict[str, Any]:
        """打开浏览器"""
        try:
            webbrowser.open(url)
            return {"success": True, "message": f"已打开浏览器: {url}"}
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_search_web(self, query: str) -> Dict[str, Any]:
        """网络搜索"""
        try:
            url = f"https://www.baidu.com/s?wd={query}"
            webbrowser.open(url)
            return {"success": True, "message": f"正在搜索: {query}"}
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_open_file(self, file_path: str) -> Dict[str, Any]:
        """打开文件"""
        try:
            os.startfile(file_path)
            return {"success": True, "message": f"已打开文件: {file_path}"}
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_open_folder(self, folder_path: str) -> Dict[str, Any]:
        """打开文件夹"""
        try:
            os.startfile(folder_path)
            return {"success": True, "message": f"已打开文件夹: {folder_path}"}
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_screenshot(self) -> Dict[str, Any]:
        """截图"""
        try:
            import pyautogui
            screenshot = pyautogui.screenshot()
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(os.path.expanduser("~"), "Pictures", f"screenshot_{timestamp}.png")
            screenshot.save(output_path)
            return {"success": True, "message": "截图成功", "result": {"path": output_path}}
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_open_application(self, app_name: str) -> Dict[str, Any]:
        """打开应用"""
        try:
            os.system(f"start {app_name}")
            return {"success": True, "message": f"正在打开应用: {app_name}"}
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_create_folder(self, folder_path: str) -> Dict[str, Any]:
        """创建文件夹"""
        try:
            os.makedirs(folder_path, exist_ok=True)
            return {"success": True, "message": f"已创建文件夹: {folder_path}"}
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_delete_file(self, file_path: str) -> Dict[str, Any]:
        """删除文件"""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                return {"success": True, "message": f"已删除文件: {file_path}"}
            else:
                return {"success": False, "message": f"文件不存在: {file_path}"}
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_copy_file(self, source_path: str, destination_path: str) -> Dict[str, Any]:
        """复制文件"""
        try:
            import shutil
            shutil.copy2(source_path, destination_path)
            return {"success": True, "message": f"已复制文件: {source_path} -> {destination_path}"}
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_generate_document(self, title: str, content: str) -> Dict[str, Any]:
        """生成Word文档"""
        try:
            from docx import Document
            doc = Document()
            doc.add_heading(title, level=1)
            doc.add_paragraph(content)
            
            output_dir = os.path.join(os.path.expanduser("~"), "Documents", "zhiy_output")
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, f"{title}.docx")
            doc.save(output_path)
            
            return {"success": True, "message": f"已生成Word文档: {output_path}", "result": {"path": output_path}}
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_generate_excel(self, data: list, headers: list = None, filename: str = "output") -> Dict[str, Any]:
        """生成Excel"""
        try:
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            
            if headers:
                ws.append(headers)
            
            for row in data:
                ws.append(row)
            
            output_dir = os.path.join(os.path.expanduser("~"), "Documents", "zhiy_output")
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, f"{filename}.xlsx")
            wb.save(output_path)
            
            return {"success": True, "message": f"已生成Excel: {output_path}", "result": {"path": output_path}}
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_generate_pptx(self, title: str, content: str) -> Dict[str, Any]:
        """生成PPT"""
        try:
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            
            title_shape = slide.shapes.title
            title_shape.text = title
            
            content_shape = slide.placeholders[1]
            content_shape.text = content
            
            output_dir = os.path.join(os.path.expanduser("~"), "Documents", "zhiy_output")
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, f"{title}.pptx")
            prs.save(output_path)
            
            return {"success": True, "message": f"已生成PPT: {output_path}", "result": {"path": output_path}}
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_get_time(self) -> Dict[str, Any]:
        """获取时间"""
        now = datetime.datetime.now()
        time_str = now.strftime("%H:%M:%S")
        return {"success": True, "message": f"当前时间: {time_str}", "result": {"time": time_str, "datetime": now}}
    
    def _execute_get_date(self) -> Dict[str, Any]:
        """获取日期"""
        now = datetime.datetime.now()
        date_str = now.strftime("%Y-%m-%d")
        weekday = now.strftime("%A")
        return {"success": True, "message": f"今天是 {date_str}，{weekday}", "result": {"date": date_str, "weekday": weekday}}
    
    def _record_tool_usage(self, tool_id: str, parameters: Dict[str, Any], result: Dict[str, Any]):
        """记录工具使用情况"""
        usage_record = {
            "tool_id": tool_id,
            "parameters": parameters,
            "result": result,
            "timestamp": datetime.datetime.now().isoformat(),
            "success": result.get('success', False)
        }
        self.tool_usage_history.append(usage_record)
        if len(self.tool_usage_history) > 1000:
            self.tool_usage_history = self.tool_usage_history[-1000:]
    
    def discover_skills(self) -> int:
        """自动发现并注册技能"""
        try:
            skills_path = Path(self.skills_dir)
            if not skills_path.exists():
                logger.warning(f"技能目录不存在: {self.skills_dir}")
                return 0
            
            discovered_count = 0
            
            for skill_file in skills_path.rglob("SKILL.md"):
                try:
                    skill_info = self._parse_skill_md(skill_file)
                    if skill_info:
                        skill_name = skill_info.get('name', skill_file.parent.name)
                        self.skills[skill_name] = skill_info
                        self.skill_registry[skill_name] = skill_info
                        discovered_count += 1
                        logger.info(f"发现技能: {skill_name}")
                except Exception as e:
                    logger.warning(f"解析技能文件失败 {skill_file}: {str(e)}")
            
            logger.info(f"发现并注册了 {discovered_count} 个技能")
            return discovered_count
            
        except Exception as e:
            logger.error(f"发现技能失败: {str(e)}")
            return 0
    
    def _parse_skill_md(self, skill_file: Path) -> Optional[Dict[str, Any]]:
        """解析 SKILL.md 文件"""
        try:
            with open(skill_file, 'r', encoding='utf-8') as f:
                content = f.read()
            
            skill_info = {
                "name": skill_file.parent.name,
                "path": str(skill_file.parent),
                "status": "available",
                "description": "",
                "category": "general",
                "parameters": {}
            }
            
            lines = content.split('\n')
            current_section = None
            
            for line in lines:
                line = line.strip()
                if line.startswith('# '):
                    skill_info['name'] = line[2:].strip()
                elif line.startswith('## '):
                    current_section = line[3:].strip().lower()
                elif current_section == '功能' and line:
                    skill_info['description'] += line + ' '
                elif current_section == '分类' and line:
                    skill_info['category'] = line
                elif current_section == '参数' and line.startswith('-'):
                    param_match = line[1:].strip()
                    if param_match:
                        skill_info['parameters'][param_match.split(':')[0].strip()] = {
                            'description': ':'.join(param_match.split(':')[1:]).strip() if ':' in param_match else ''
                        }
            
            return skill_info
            
        except Exception as e:
            logger.error(f"解析技能文件失败: {str(e)}")
            return None
    
    def get_skill(self, skill_name: str) -> Optional[Dict[str, Any]]:
        """获取技能"""
        return self.skills.get(skill_name)
    
    def list_skills(self, category: str = None) -> List[Dict[str, Any]]:
        """列出所有技能"""
        if category:
            return [s for s in self.skill_registry.values() if s.get('category') == category]
        return list(self.skill_registry.values())
    
    def execute_skill(self, skill_name: str, **kwargs) -> Dict[str, Any]:
        """执行技能"""
        try:
            skill = self.get_skill(skill_name)
            if not skill:
                return {"success": False, "message": f"技能 {skill_name} 不存在"}
            
            skill_path = skill.get('path', '')
            main_script = os.path.join(skill_path, 'main.py')
            
            if os.path.exists(main_script):
                args = ' '.join([f'--{k} "{v}"' for k, v in kwargs.items()])
                result = subprocess.run(
                    ['python', main_script] + (args.split() if args else []),
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                
                if result.returncode == 0:
                    return {"success": True, "message": f"技能 {skill_name} 执行成功", "output": result.stdout}
                else:
                    return {"success": False, "message": f"技能执行失败: {result.stderr}"}
            else:
                return {"success": False, "message": f"技能脚本不存在: {main_script}"}
                
        except Exception as e:
            logger.error(f"执行技能 {skill_name} 失败: {str(e)}")
            return {"success": False, "message": f"技能执行失败: {str(e)}"}
    
    def get_statistics(self) -> Dict[str, Any]:
        """获取统计信息"""
        return {
            "tools": {
                "local": len(self.tools['local']),
                "external": len(self.tools['external'])
            },
            "skills": {
                "total": len(self.skills),
                "categories": len(set(s.get('category', 'general') for s in self.skills.values()))
            },
            "usage_history": len(self.tool_usage_history)
        }


openclaw_core = OpenClawCore()
