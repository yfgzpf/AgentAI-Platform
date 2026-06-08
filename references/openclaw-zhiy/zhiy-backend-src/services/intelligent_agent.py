"""
智 Y.Ai 智能体编排器
实现真正的 OpenClaw 自动化能力：
1. 意图识别 → 自动匹配技能
2. 参数缺失 → 智能追问
3. 任务分解 → 自动执行
4. 流式输出 → 实时反馈
"""

import os
import sys
import json
import re
import asyncio
from typing import Dict, Any, List, Optional, Tuple
from dataclasses import dataclass, field
from enum import Enum
from datetime import datetime
from pathlib import Path

# 添加当前目录到路径
CURRENT_DIR = Path(__file__).parent
if str(CURRENT_DIR) not in sys.path:
    sys.path.insert(0, str(CURRENT_DIR))

try:
    from loguru import logger
except ImportError:
    import logging
    logger = logging.getLogger(__name__)


class IntentType(Enum):
    """意图类型"""
    BROWSER_AUTOMATION = "browser_automation"
    DESKTOP_AUTOMATION = "desktop_automation"
    DOCUMENT_GENERATION = "document_generation"
    IMAGE_GENERATION = "image_generation"
    VIDEO_GENERATION = "video_generation"
    FILE_OPERATION = "file_operation"
    COMMUNICATION = "communication"
    CODE_EXECUTION = "code_execution"
    QUOTATION_GENERATION = "quotation_generation"
    CONTRACT_GENERATION = "contract_generation"
    GENERAL_CHAT = "general_chat"
    UNKNOWN = "unknown"


@dataclass
class ParsedIntent:
    """解析后的意图"""
    intent_type: IntentType
    confidence: float
    entities: Dict[str, Any]
    missing_params: List[str]
    suggested_skill: Optional[str]
    action_sequence: List[Dict[str, Any]]


class IntelligentAgentOrchestrator:
    """智能体编排器 - 实现 OpenClaw 核心自动化能力"""
    
    def __init__(self):
        self.intent_patterns = self._build_intent_patterns()
        self.skill_registry = self._build_skill_registry()
        self.conversation_context: Dict[str, Any] = {}
        logger.info("智能体编排器初始化完成")
    
    def _build_intent_patterns(self) -> Dict[IntentType, List[Dict[str, Any]]]:
        """构建意图模式库"""
        return {
            IntentType.BROWSER_AUTOMATION: [
                {
                    "patterns": [r"打开(.+?)(网页|网站|浏览器)", r"访问(.+)", r"去(.+?)看"],
                    "params": {"url": 1},
                    "skill": "browser_automation",
                    "action": "open"
                },
                {
                    "patterns": [r"搜索(.+)", r"查找(.+)", r"搜一下(.+)"],
                    "params": {"query": 1},
                    "skill": "browser_automation",
                    "action": "search"
                },
                {
                    "patterns": [r"截图", r"截屏", r"抓图"],
                    "params": {},
                    "skill": "browser_automation",
                    "action": "screenshot"
                }
            ],
            IntentType.DESKTOP_AUTOMATION: [
                {
                    "patterns": [r"打开(.+?)(应用|软件|程序|APP)", r"启动(.+)", r"运行(.+)"],
                    "params": {"app_name": 1},
                    "skill": "desktop_automation",
                    "action": "open_app"
                },
                {
                    "patterns": [r"截个图", r"截屏", r"屏幕截图"],
                    "params": {},
                    "skill": "desktop_automation",
                    "action": "screenshot"
                }
            ],
            IntentType.DOCUMENT_GENERATION: [
                {
                    "patterns": [r"生成(.+?)(文档|Word|报告)", r"创建(.+?)(文档|Word)", r"写(.+?)(文档|报告)"],
                    "params": {"title": 1, "content": None},
                    "skill": "document_generation",
                    "action": "generate"
                },
                {
                    "patterns": [r"制作(.+?)PPT", r"生成(.+?)演示文稿", r"创建(.+?)PPT"],
                    "params": {"title": 1, "content": None},
                    "skill": "document_generation",
                    "action": "generate_pptx"
                },
                {
                    "patterns": [r"生成(.+?)表格", r"创建(.+?)Excel", r"制作(.+?)报表"],
                    "params": {"title": 1, "data": None},
                    "skill": "document_generation",
                    "action": "generate_excel"
                }
            ],
            IntentType.IMAGE_GENERATION: [
                {
                    "patterns": [r"生成(.+?)(图片|图像|画作)", r"画(.+?)(图|画)", r"AI画(.+)", r"创建(.+?)图片"],
                    "params": {"prompt": 1},
                    "skill": "image_generation",
                    "action": "generate"
                }
            ],
            IntentType.VIDEO_GENERATION: [
                {
                    "patterns": [r"生成(.+?)(视频|影片)", r"制作(.+?)视频", r"AI视频(.+)", r"用seed生成(.+)"],
                    "params": {"prompt": 1},
                    "skill": "video_generation",
                    "action": "generate"
                }
            ],
            IntentType.QUOTATION_GENERATION: [
                {
                    "patterns": [r"生成(.+?)报价", r"制作(.+?)报价单", r"装修报价", r"建材报价"],
                    "params": {"customerName": None, "area": None, "style": None},
                    "skill": "quotation_generation",
                    "action": "generate"
                }
            ],
            IntentType.CONTRACT_GENERATION: [
                {
                    "patterns": [r"生成(.+?)合同", r"制作(.+?)合同", r"写(.+?)合同"],
                    "params": {"customerName": None, "area": None, "style": None},
                    "skill": "contract_generation",
                    "action": "generate"
                }
            ],
            IntentType.FILE_OPERATION: [
                {
                    "patterns": [r"创建(.+?)文件", r"新建(.+?)文件"],
                    "params": {"path": 1},
                    "skill": "file_operation",
                    "action": "create"
                },
                {
                    "patterns": [r"读取(.+?)文件", r"打开(.+?)文件"],
                    "params": {"path": 1},
                    "skill": "file_operation",
                    "action": "read"
                },
                {
                    "patterns": [r"删除(.+?)文件", r"移除(.+?)文件"],
                    "params": {"path": 1},
                    "skill": "file_operation",
                    "action": "delete"
                }
            ],
            IntentType.COMMUNICATION: [
                {
                    "patterns": [r"发(.+?)微信", r"发送(.+?)微信", r"微信(.+?)说"],
                    "params": {"contact": None, "message": None},
                    "skill": "communication",
                    "action": "wechat"
                },
                {
                    "patterns": [r"发(.+?)邮件", r"发送(.+?)邮件", r"邮件(.+?)说"],
                    "params": {"recipient": None, "subject": None, "content": None},
                    "skill": "communication",
                    "action": "email"
                }
            ],
            IntentType.CODE_EXECUTION: [
                {
                    "patterns": [r"运行(.+?)代码", r"执行(.+?)代码", r"跑(.+?)脚本"],
                    "params": {"code": None, "language": None},
                    "skill": "code_execution",
                    "action": "execute"
                }
            ]
        }
    
    def _build_skill_registry(self) -> Dict[str, Dict[str, Any]]:
        """构建技能注册表"""
        return {
            "browser_automation": {
                "name": "浏览器自动化",
                "description": "打开网页、搜索、截图等",
                "executor": "_execute_browser_automation",
                "params_schema": {
                    "url": {"type": "string", "required": False, "description": "网址"},
                    "query": {"type": "string", "required": False, "description": "搜索关键词"},
                    "action": {"type": "string", "required": True, "description": "操作类型"}
                }
            },
            "desktop_automation": {
                "name": "桌面自动化",
                "description": "打开应用、截图等",
                "executor": "_execute_desktop_automation",
                "params_schema": {
                    "app_name": {"type": "string", "required": False, "description": "应用名称"},
                    "action": {"type": "string", "required": True, "description": "操作类型"}
                }
            },
            "document_generation": {
                "name": "文档生成",
                "description": "生成Word、Excel、PPT文档",
                "executor": "_execute_document_generation",
                "params_schema": {
                    "title": {"type": "string", "required": True, "description": "文档标题"},
                    "content": {"type": "string", "required": False, "description": "文档内容"},
                    "doc_type": {"type": "string", "required": True, "description": "文档类型"}
                }
            },
            "image_generation": {
                "name": "图像生成",
                "description": "AI生成图片",
                "executor": "_execute_image_generation",
                "params_schema": {
                    "prompt": {"type": "string", "required": True, "description": "图像描述"}
                }
            },
            "video_generation": {
                "name": "视频生成",
                "description": "AI生成视频",
                "executor": "_execute_video_generation",
                "params_schema": {
                    "prompt": {"type": "string", "required": True, "description": "视频描述"}
                }
            },
            "quotation_generation": {
                "name": "报价单生成",
                "description": "生成装修报价单",
                "executor": "_execute_quotation_generation",
                "params_schema": {
                    "customerName": {"type": "string", "required": True, "description": "客户姓名"},
                    "area": {"type": "number", "required": True, "description": "装修面积"},
                    "style": {"type": "string", "required": True, "description": "装修风格"}
                }
            },
            "contract_generation": {
                "name": "合同生成",
                "description": "生成装修合同",
                "executor": "_execute_contract_generation",
                "params_schema": {
                    "customerName": {"type": "string", "required": True, "description": "客户姓名"},
                    "area": {"type": "number", "required": True, "description": "装修面积"},
                    "style": {"type": "string", "required": True, "description": "装修风格"}
                }
            },
            "file_operation": {
                "name": "文件操作",
                "description": "创建、读取、删除文件",
                "executor": "_execute_file_operation",
                "params_schema": {
                    "path": {"type": "string", "required": True, "description": "文件路径"},
                    "action": {"type": "string", "required": True, "description": "操作类型"}
                }
            },
            "communication": {
                "name": "通信",
                "description": "发送微信、邮件",
                "executor": "_execute_communication",
                "params_schema": {
                    "platform": {"type": "string", "required": True, "description": "平台"},
                    "recipient": {"type": "string", "required": True, "description": "接收者"},
                    "message": {"type": "string", "required": True, "description": "消息内容"}
                }
            },
            "code_execution": {
                "name": "代码执行",
                "description": "运行代码",
                "executor": "_execute_code",
                "params_schema": {
                    "code": {"type": "string", "required": True, "description": "代码内容"},
                    "language": {"type": "string", "required": True, "description": "编程语言"}
                }
            }
        }
    
    def parse_intent(self, user_input: str, context: Dict[str, Any] = None) -> ParsedIntent:
        """
        解析用户意图
        
        Args:
            user_input: 用户输入
            context: 上下文
        
        Returns:
            解析后的意图
        """
        user_input_lower = user_input.lower()
        
        # 遍历所有意图模式
        for intent_type, patterns_list in self.intent_patterns.items():
            for pattern_info in patterns_list:
                for pattern in pattern_info["patterns"]:
                    match = re.search(pattern, user_input)
                    if match:
                        # 提取实体
                        entities = {}
                        for param_name, group_idx in pattern_info["params"].items():
                            if group_idx is not None and match.groups():
                                try:
                                    entities[param_name] = match.group(group_idx)
                                except IndexError:
                                    pass
                        
                        entities["action"] = pattern_info["action"]
                        
                        # 检查缺失参数
                        skill_name = pattern_info["skill"]
                        skill_info = self.skill_registry.get(skill_name, {})
                        params_schema = skill_info.get("params_schema", {})
                        
                        missing_params = []
                        for param_name, param_info in params_schema.items():
                            if param_info.get("required", False) and param_name not in entities:
                                missing_params.append(param_name)
                        
                        # 构建动作序列
                        action_sequence = [
                            {
                                "step": 1,
                                "action": "analyze",
                                "description": f"分析用户意图: {intent_type.value}",
                                "status": "completed"
                            },
                            {
                                "step": 2,
                                "action": "prepare",
                                "description": f"准备执行技能: {skill_name}",
                                "status": "pending"
                            },
                            {
                                "step": 3,
                                "action": "execute",
                                "description": f"执行技能: {skill_name}",
                                "skill": skill_name,
                                "params": entities,
                                "status": "pending"
                            }
                        ]
                        
                        return ParsedIntent(
                            intent_type=intent_type,
                            confidence=0.85,
                            entities=entities,
                            missing_params=missing_params,
                            suggested_skill=skill_name,
                            action_sequence=action_sequence
                        )
        
        # 未识别到特定意图，返回通用聊天
        return ParsedIntent(
            intent_type=IntentType.GENERAL_CHAT,
            confidence=0.5,
            entities={},
            missing_params=[],
            suggested_skill=None,
            action_sequence=[
                {
                    "step": 1,
                    "action": "chat",
                    "description": "进行对话",
                    "status": "pending"
                }
            ]
        )
    
    def generate_questions(self, missing_params: List[str], skill_name: str) -> List[str]:
        """生成追问问题"""
        skill_info = self.skill_registry.get(skill_name, {})
        params_schema = skill_info.get("params_schema", {})
        
        questions = []
        for param_name in missing_params:
            param_info = params_schema.get(param_name, {})
            description = param_info.get("description", param_name)
            param_type = param_info.get("type", "string")
            
            if param_type == "number":
                questions.append(f"请提供{description}（数字）：")
            elif param_type == "choice":
                options = param_info.get("options", [])
                questions.append(f"请选择{description}：{', '.join(options)}")
            else:
                questions.append(f"请提供{description}：")
        
        return questions
    
    async def execute_skill(self, skill_name: str, params: Dict[str, Any], stream_callback=None) -> Dict[str, Any]:
        """
        执行技能
        
        Args:
            skill_name: 技能名称
            params: 参数
            stream_callback: 流式回调函数
        
        Returns:
            执行结果
        """
        skill_info = self.skill_registry.get(skill_name)
        if not skill_info:
            return {"success": False, "message": f"技能 {skill_name} 不存在"}
        
        executor_name = skill_info.get("executor")
        if not executor_name:
            return {"success": False, "message": f"技能 {skill_name} 没有定义执行器"}
        
        # 流式输出开始
        if stream_callback:
            await stream_callback({
                "type": "skill_start",
                "skill": skill_name,
                "message": f"正在执行技能: {skill_info['name']}..."
            })
        
        # 执行技能
        try:
            executor = getattr(self, executor_name, None)
            if executor:
                result = await executor(params, stream_callback)
            else:
                result = {"success": False, "message": f"执行器 {executor_name} 未实现"}
        except Exception as e:
            result = {"success": False, "message": f"执行失败: {str(e)}"}
        
        # 流式输出结束
        if stream_callback:
            await stream_callback({
                "type": "skill_end",
                "skill": skill_name,
                "result": result
            })
        
        return result
    
    async def _execute_browser_automation(self, params: Dict[str, Any], stream_callback=None) -> Dict[str, Any]:
        """执行浏览器自动化"""
        import webbrowser
        
        action = params.get("action", "open")
        
        if stream_callback:
            await stream_callback({
                "type": "thinking",
                "message": f"正在执行浏览器操作: {action}..."
            })
        
        if action == "open":
            url = params.get("url", "https://www.baidu.com")
            if not url.startswith("http"):
                url = f"https://{url}"
            webbrowser.open(url)
            return {"success": True, "message": f"已打开浏览器: {url}"}
        
        elif action == "search":
            query = params.get("query", "")
            url = f"https://www.baidu.com/s?wd={query}"
            webbrowser.open(url)
            return {"success": True, "message": f"正在搜索: {query}"}
        
        elif action == "screenshot":
            try:
                import pyautogui
                screenshot = pyautogui.screenshot()
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_path = os.path.join(os.path.expanduser("~"), "Pictures", f"screenshot_{timestamp}.png")
                screenshot.save(output_path)
                return {"success": True, "message": f"截图已保存: {output_path}", "path": output_path}
            except ImportError:
                return {"success": False, "message": "pyautogui 未安装"}
        
        return {"success": False, "message": f"未知操作: {action}"}
    
    async def _execute_desktop_automation(self, params: Dict[str, Any], stream_callback=None) -> Dict[str, Any]:
        """执行桌面自动化"""
        action = params.get("action", "open_app")
        
        if stream_callback:
            await stream_callback({
                "type": "thinking",
                "message": f"正在执行桌面操作: {action}..."
            })
        
        if action == "open_app":
            app_name = params.get("app_name", "")
            if not app_name:
                return {"success": False, "message": "请指定要打开的应用名称"}
            
            try:
                os.system(f"start {app_name}")
                return {"success": True, "message": f"正在打开应用: {app_name}"}
            except Exception as e:
                return {"success": False, "message": str(e)}
        
        elif action == "screenshot":
            try:
                import pyautogui
                screenshot = pyautogui.screenshot()
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_path = os.path.join(os.path.expanduser("~"), "Pictures", f"screenshot_{timestamp}.png")
                screenshot.save(output_path)
                return {"success": True, "message": f"截图已保存: {output_path}", "path": output_path}
            except ImportError:
                return {"success": False, "message": "pyautogui 未安装"}
        
        return {"success": False, "message": f"未知操作: {action}"}
    
    async def _execute_document_generation(self, params: Dict[str, Any], stream_callback=None) -> Dict[str, Any]:
        """执行文档生成"""
        doc_type = params.get("doc_type", "word")
        title = params.get("title", "未命名文档")
        content = params.get("content", "")
        
        if stream_callback:
            await stream_callback({
                "type": "thinking",
                "message": f"正在生成{doc_type}文档..."
            })
        
        output_dir = os.path.join(os.path.expanduser("~"), "Documents", "zhiy_output")
        os.makedirs(output_dir, exist_ok=True)
        
        try:
            if doc_type in ["word", "docx"]:
                from docx import Document
                doc = Document()
                doc.add_heading(title, level=1)
                doc.add_paragraph(content)
                output_path = os.path.join(output_dir, f"{title}.docx")
                doc.save(output_path)
                return {"success": True, "message": f"Word文档已生成: {output_path}", "path": output_path}
            
            elif doc_type in ["excel", "xlsx"]:
                import openpyxl
                wb = openpyxl.Workbook()
                ws = wb.active
                ws.title = title
                ws['A1'] = title
                output_path = os.path.join(output_dir, f"{title}.xlsx")
                wb.save(output_path)
                return {"success": True, "message": f"Excel文档已生成: {output_path}", "path": output_path}
            
            elif doc_type in ["ppt", "pptx"]:
                from pptx import Presentation
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = title
                output_path = os.path.join(output_dir, f"{title}.pptx")
                prs.save(output_path)
                return {"success": True, "message": f"PPT已生成: {output_path}", "path": output_path}
            
        except ImportError as e:
            return {"success": False, "message": f"缺少依赖: {str(e)}"}
        except Exception as e:
            return {"success": False, "message": str(e)}
        
        return {"success": False, "message": f"不支持的文档类型: {doc_type}"}
    
    async def _execute_image_generation(self, params: Dict[str, Any], stream_callback=None) -> Dict[str, Any]:
        """执行图像生成"""
        prompt = params.get("prompt", "")
        
        if stream_callback:
            await stream_callback({
                "type": "thinking",
                "message": f"正在生成图像: {prompt}..."
            })
        
        # 调用图像生成 API
        return {
            "success": False,
            "message": "图像生成需要配置 API 密钥。请设置 DEEPSEEK_API_KEY 或 QIANWEN_API_KEY 环境变量。",
            "prompt": prompt
        }
    
    async def _execute_video_generation(self, params: Dict[str, Any], stream_callback=None) -> Dict[str, Any]:
        """执行视频生成"""
        prompt = params.get("prompt", "")
        
        if stream_callback:
            await stream_callback({
                "type": "thinking",
                "message": f"正在生成视频: {prompt}..."
            })
        
        # 调用视频生成 API
        return {
            "success": False,
            "message": "视频生成需要配置豆包 API 密钥。请设置 DOUBAO_API_KEY 环境变量。",
            "prompt": prompt
        }
    
    async def _execute_quotation_generation(self, params: Dict[str, Any], stream_callback=None) -> Dict[str, Any]:
        """执行报价单生成"""
        customer_name = params.get("customerName", "客户")
        area = params.get("area", 100)
        style = params.get("style", "现代简约")
        
        if stream_callback:
            await stream_callback({
                "type": "thinking",
                "message": f"正在生成报价单: {customer_name}, {area}㎡, {style}..."
            })
        
        # 调用报价单生成技能
        skill_path = os.path.join(os.path.dirname(__file__), "..", "..", "skills", "projects", "construction", "quotation", "main.py")
        if os.path.exists(skill_path):
            import subprocess
            result = subprocess.run(
                ["python", skill_path, "--customerName", customer_name, "--area", str(area), "--style", style],
                capture_output=True,
                text=True,
                timeout=60
            )
            if result.returncode == 0:
                return {"success": True, "message": "报价单生成成功", "output": result.stdout}
            else:
                return {"success": False, "message": result.stderr}
        
        return {"success": False, "message": "报价单生成技能未找到"}
    
    async def _execute_contract_generation(self, params: Dict[str, Any], stream_callback=None) -> Dict[str, Any]:
        """执行合同生成"""
        customer_name = params.get("customerName", "客户")
        area = params.get("area", 100)
        style = params.get("style", "现代简约")
        
        if stream_callback:
            await stream_callback({
                "type": "thinking",
                "message": f"正在生成合同: {customer_name}, {area}㎡, {style}..."
            })
        
        return {"success": True, "message": f"合同生成功能开发中..."}
    
    async def _execute_file_operation(self, params: Dict[str, Any], stream_callback=None) -> Dict[str, Any]:
        """执行文件操作"""
        action = params.get("action", "create")
        path = params.get("path", "")
        
        if stream_callback:
            await stream_callback({
                "type": "thinking",
                "message": f"正在执行文件操作: {action} {path}..."
            })
        
        try:
            if action == "create":
                os.makedirs(os.path.dirname(path), exist_ok=True)
                with open(path, 'w', encoding='utf-8') as f:
                    f.write("")
                return {"success": True, "message": f"文件已创建: {path}"}
            
            elif action == "read":
                if os.path.exists(path):
                    with open(path, 'r', encoding='utf-8') as f:
                        content = f.read()
                    return {"success": True, "message": f"文件内容:", "content": content}
                return {"success": False, "message": f"文件不存在: {path}"}
            
            elif action == "delete":
                if os.path.exists(path):
                    os.remove(path)
                    return {"success": True, "message": f"文件已删除: {path}"}
                return {"success": False, "message": f"文件不存在: {path}"}
            
        except Exception as e:
            return {"success": False, "message": str(e)}
        
        return {"success": False, "message": f"未知操作: {action}"}
    
    async def _execute_communication(self, params: Dict[str, Any], stream_callback=None) -> Dict[str, Any]:
        """执行通信操作"""
        platform = params.get("platform", "wechat")
        recipient = params.get("recipient", "")
        message = params.get("message", "")
        
        if stream_callback:
            await stream_callback({
                "type": "thinking",
                "message": f"正在发送{platform}消息..."
            })
        
        return {
            "success": False,
            "message": f"{platform}消息发送需要配置相关服务。"
        }
    
    async def _execute_code(self, params: Dict[str, Any], stream_callback=None) -> Dict[str, Any]:
        """执行代码"""
        code = params.get("code", "")
        language = params.get("language", "python")
        
        if stream_callback:
            await stream_callback({
                "type": "thinking",
                "message": f"正在执行{language}代码..."
            })
        
        if language == "python":
            try:
                exec_globals = {}
                exec(code, exec_globals)
                return {"success": True, "message": "代码执行成功", "result": exec_globals}
            except Exception as e:
                return {"success": False, "message": f"执行错误: {str(e)}"}
        
        return {"success": False, "message": f"不支持的语言: {language}"}
    
    async def process_user_input(
        self,
        user_input: str,
        context: Dict[str, Any] = None,
        stream_callback=None
    ) -> Dict[str, Any]:
        """
        处理用户输入 - 核心自动化入口
        
        Args:
            user_input: 用户输入
            context: 上下文
            stream_callback: 流式回调
        
        Returns:
            处理结果
        """
        # 1. 解析意图
        if stream_callback:
            await stream_callback({
                "type": "thinking",
                "message": "正在分析您的请求..."
            })
        
        parsed_intent = self.parse_intent(user_input, context)
        
        # 2. 检查是否需要追问
        if parsed_intent.missing_params:
            questions = self.generate_questions(parsed_intent.missing_params, parsed_intent.suggested_skill)
            return {
                "success": True,
                "type": "question",
                "message": f"我理解您想要**{parsed_intent.intent_type.value}**。\n\n请提供以下信息：",
                "questions": questions,
                "missing_params": parsed_intent.missing_params,
                "suggested_skill": parsed_intent.suggested_skill
            }
        
        # 3. 执行技能
        if parsed_intent.suggested_skill:
            result = await self.execute_skill(
                parsed_intent.suggested_skill,
                parsed_intent.entities,
                stream_callback
            )
            return {
                "success": result.get("success", False),
                "type": "result",
                "message": result.get("message", ""),
                "result": result
            }
        
        # 4. 通用对话
        return {
            "success": True,
            "type": "chat",
            "message": f"我理解您的请求。请问您需要什么帮助？"
        }


# 全局实例
intelligent_agent = IntelligentAgentOrchestrator()
