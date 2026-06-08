"""
工具执行器
为工具注册表中的每个工具提供实际的执行函数
"""

import os
import sys
import datetime
import webbrowser
import subprocess
from typing import Dict, Any
from pathlib import Path

# 添加当前目录到路径
CURRENT_DIR = Path(__file__).parent
if str(CURRENT_DIR) not in sys.path:
    sys.path.insert(0, str(CURRENT_DIR))

try:
    from loguru import logger
except ImportError:
    import logging
    logger = logging.getLogger(__name__)

# 尝试导入服务，如果失败则使用占位符
try:
    from gui_automation_service import gui_automation_service
    GUI_SERVICE_AVAILABLE = True
except ImportError:
    gui_automation_service = None
    GUI_SERVICE_AVAILABLE = False

try:
    from desktop_automation_service import desktop_automation_service
    DESKTOP_SERVICE_AVAILABLE = True
except ImportError:
    desktop_automation_service = None
    DESKTOP_SERVICE_AVAILABLE = False


class ToolExecutor:
    """工具执行器"""
    
    def __init__(self):
        self.gui_service = gui_automation_service
        self.desktop_service = desktop_automation_service
        logger.info("工具执行器初始化完成")
    
    def execute(self, tool_id: str, **kwargs) -> Dict[str, Any]:
        """
        执行工具
        
        Args:
            tool_id: 工具ID
            **kwargs: 工具参数
        
        Returns:
            执行结果
        """
        try:
            logger.info(f"执行工具: {tool_id}, 参数: {kwargs}")
            
            # 获取对应的执行方法
            method_name = f"_execute_{tool_id}"
            if hasattr(self, method_name):
                method = getattr(self, method_name)
                return method(**kwargs)
            else:
                return {
                    "success": False,
                    "message": f"未找到工具 {tool_id} 的执行方法"
                }
        
        except Exception as e:
            logger.error(f"执行工具 {tool_id} 失败: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            return {
                "success": False,
                "message": f"执行失败: {str(e)}"
            }
    
    # ===== 桌面自动化工具 =====
    
    def _execute_play_music(self, song_name: str = None, search_query: str = None) -> Dict[str, Any]:
        """播放音乐"""
        try:
            query = song_name or search_query or "音乐"
            result = self.desktop_service.play_music(query)
            return {
                "success": True,
                "message": f"正在播放: {query}",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_open_browser(self, url: str = "https://www.baidu.com") -> Dict[str, Any]:
        """打开浏览器"""
        try:
            webbrowser.open(url)
            return {
                "success": True,
                "message": f"已打开浏览器: {url}"
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_search_web(self, query: str) -> Dict[str, Any]:
        """网络搜索"""
        try:
            url = f"https://www.baidu.com/s?wd={query}"
            webbrowser.open(url)
            return {
                "success": True,
                "message": f"正在搜索: {query}"
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_open_file(self, file_path: str) -> Dict[str, Any]:
        """打开文件"""
        try:
            os.startfile(file_path)
            return {
                "success": True,
                "message": f"已打开文件: {file_path}"
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_open_folder(self, folder_path: str) -> Dict[str, Any]:
        """打开文件夹"""
        try:
            os.startfile(folder_path)
            return {
                "success": True,
                "message": f"已打开文件夹: {folder_path}"
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_screenshot(self) -> Dict[str, Any]:
        """截图"""
        try:
            result = self.gui_service.screenshot()
            return {
                "success": True,
                "message": "截图成功",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_send_notification(self, title: str, message: str) -> Dict[str, Any]:
        """发送通知"""
        try:
            result = self.desktop_service.send_notification(title, message)
            return {
                "success": True,
                "message": "通知已发送",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_open_application(self, app_name: str, app_type: str = None) -> Dict[str, Any]:
        """打开应用"""
        try:
            result = self.desktop_service.open_application(app_name)
            return {
                "success": True,
                "message": f"正在打开应用: {app_name}",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    # ===== GUI自动化工具 =====
    
    def _execute_mouse_move(self, x: int, y: int) -> Dict[str, Any]:
        """移动鼠标"""
        try:
            result = self.gui_service.mouse_move(x, y)
            return {
                "success": True,
                "message": f"鼠标已移动到 ({x}, {y})",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_mouse_click(self, x: int = None, y: int = None, button: str = "left", clicks: int = 1) -> Dict[str, Any]:
        """点击鼠标"""
        try:
            if x is not None and y is not None:
                self.gui_service.mouse_move(x, y)
            result = self.gui_service.mouse_click(button=button, clicks=clicks)
            return {
                "success": True,
                "message": f"鼠标已点击 ({button}键, {clicks}次)",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_mouse_double_click(self, x: int = None, y: int = None) -> Dict[str, Any]:
        """双击鼠标"""
        try:
            if x is not None and y is not None:
                self.gui_service.mouse_move(x, y)
            result = self.gui_service.mouse_double_click()
            return {
                "success": True,
                "message": "鼠标已双击",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_mouse_right_click(self, x: int = None, y: int = None) -> Dict[str, Any]:
        """右键点击"""
        try:
            if x is not None and y is not None:
                self.gui_service.mouse_move(x, y)
            result = self.gui_service.mouse_right_click()
            return {
                "success": True,
                "message": "鼠标已右键点击",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_mouse_scroll(self, clicks: int = 1) -> Dict[str, Any]:
        """滚动鼠标"""
        try:
            result = self.gui_service.mouse_scroll(clicks)
            return {
                "success": True,
                "message": f"鼠标已滚动 {clicks} 次",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_mouse_drag(self, x: int, y: int) -> Dict[str, Any]:
        """拖拽鼠标"""
        try:
            result = self.gui_service.mouse_drag(x, y)
            return {
                "success": True,
                "message": f"鼠标已拖拽到 ({x}, {y})",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_type_text(self, text: str) -> Dict[str, Any]:
        """输入文本"""
        try:
            result = self.gui_service.type_text(text)
            return {
                "success": True,
                "message": f"已输入文本: {text}",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_press_key(self, key: str, presses: int = 1) -> Dict[str, Any]:
        """按下按键"""
        try:
            result = self.gui_service.press_key(key, presses)
            return {
                "success": True,
                "message": f"已按下按键: {key} ({presses}次)",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_hotkey(self, keys: list) -> Dict[str, Any]:
        """组合键"""
        try:
            result = self.gui_service.hotkey(*keys)
            return {
                "success": True,
                "message": f"已按下组合键: {keys}",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_copy_text(self, text: str) -> Dict[str, Any]:
        """复制文本"""
        try:
            result = self.gui_service.copy_text(text)
            return {
                "success": True,
                "message": f"已复制文本到剪贴板",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_paste_text(self) -> Dict[str, Any]:
        """粘贴文本"""
        try:
            result = self.gui_service.paste_text()
            return {
                "success": True,
                "message": "已从剪贴板粘贴文本",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_get_mouse_position(self) -> Dict[str, Any]:
        """获取鼠标位置"""
        try:
            result = self.gui_service.get_mouse_position()
            return {
                "success": True,
                "message": f"鼠标位置: ({result['x']}, {result['y']})",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    # ===== 系统工具 =====
    
    def _execute_get_time(self) -> Dict[str, Any]:
        """获取时间"""
        try:
            now = datetime.datetime.now()
            time_str = now.strftime("%H:%M:%S")
            return {
                "success": True,
                "message": f"当前时间: {time_str}",
                "result": {"time": time_str, "datetime": now}
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_get_date(self) -> Dict[str, Any]:
        """获取日期"""
        try:
            now = datetime.datetime.now()
            date_str = now.strftime("%Y-%m-%d")
            weekday = now.strftime("%A")
            return {
                "success": True,
                "message": f"今天是 {date_str}，{weekday}",
                "result": {"date": date_str, "weekday": weekday, "datetime": now}
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_create_folder(self, folder_path: str) -> Dict[str, Any]:
        """创建文件夹"""
        try:
            os.makedirs(folder_path, exist_ok=True)
            return {
                "success": True,
                "message": f"已创建文件夹: {folder_path}"
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_delete_file(self, file_path: str) -> Dict[str, Any]:
        """删除文件"""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                return {
                    "success": True,
                    "message": f"已删除文件: {file_path}"
                }
            else:
                return {
                    "success": False,
                    "message": f"文件不存在: {file_path}"
                }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_copy_file(self, source_path: str, destination_path: str) -> Dict[str, Any]:
        """复制文件"""
        try:
            import shutil
            shutil.copy2(source_path, destination_path)
            return {
                "success": True,
                "message": f"已复制文件: {source_path} -> {destination_path}"
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_move_file(self, source_path: str, destination_path: str) -> Dict[str, Any]:
        """移动文件"""
        try:
            import shutil
            shutil.move(source_path, destination_path)
            return {
                "success": True,
                "message": f"已移动文件: {source_path} -> {destination_path}"
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_list_files(self, folder_path: str) -> Dict[str, Any]:
        """列出文件"""
        try:
            if os.path.exists(folder_path):
                files = os.listdir(folder_path)
                return {
                    "success": True,
                    "message": f"文件夹 {folder_path} 中的文件:",
                    "result": {"files": files, "count": len(files)}
                }
            else:
                return {
                    "success": False,
                    "message": f"文件夹不存在: {folder_path}"
                }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    # ===== 文本处理工具 =====
    
    def _execute_generate_document(self, title: str, content: str) -> Dict[str, Any]:
        """生成Word文档"""
        try:
            from docx import Document
            doc = Document()
            doc.add_heading(title, level=1)
            doc.add_paragraph(content)
            
            output_path = f"F:\\AI系统开发原始框架\\output\\{title}.docx"
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            doc.save(output_path)
            
            return {
                "success": True,
                "message": f"已生成Word文档: {output_path}",
                "result": {"path": output_path}
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_generate_pdf(self, title: str, content: str) -> Dict[str, Any]:
        """生成PDF文档"""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            from reportlab.lib.units import inch
            
            output_path = f"F:\\AI系统开发原始框架\\output\\{title}.pdf"
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            c = canvas.Canvas(output_path, pagesize=letter)
            c.setFont("Helvetica-Bold", 16)
            c.drawString(100, 750, title)
            c.setFont("Helvetica", 12)
            
            y = 700
            for line in content.split('\n'):
                c.drawString(100, y, line)
                y -= 20
            
            c.save()
            
            return {
                "success": True,
                "message": f"已生成PDF文档: {output_path}",
                "result": {"path": output_path}
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_generate_pptx(self, title: str, content: str) -> Dict[str, Any]:
        """生成PPT"""
        try:
            from pptx import Presentation
            
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            
            title_shape = slide.shapes.title
            title_shape.text = title
            
            content_shape = slide.placeholders[1]
            content_shape.text = content
            
            output_path = f"F:\\AI系统开发原始框架\\output\\{title}.pptx"
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            prs.save(output_path)
            
            return {
                "success": True,
                "message": f"已生成PPT: {output_path}",
                "result": {"path": output_path}
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    # ===== 其他工具（简化实现）=====
    
    def _execute_generate_image(self, prompt: str, style: str = None) -> Dict[str, Any]:
        """生成图像"""
        return {
            "success": False,
            "message": "图像生成功能需要外部API支持，暂未实现"
        }
    
    def _execute_edit_image(self, image_path: str, operations: list) -> Dict[str, Any]:
        """编辑图像"""
        return {
            "success": False,
            "message": "图像编辑功能需要外部API支持，暂未实现"
        }
    
    def _execute_generate_quotation(self, project_info: dict, items: list) -> Dict[str, Any]:
        """生成报价单"""
        try:
            from docx import Document
            doc = Document()
            doc.add_heading("装修报价单", level=1)
            
            doc.add_heading("项目信息", level=2)
            for key, value in project_info.items():
                doc.add_paragraph(f"{key}: {value}")
            
            doc.add_heading("报价明细", level=2)
            table = doc.add_table(rows=1, cols=3)
            table.style = 'Table Grid'
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = "项目"
            hdr_cells[1].text = "数量"
            hdr_cells[2].text = "价格"
            
            for item in items:
                row_cells = table.add_row().cells
                row_cells[0].text = item.get("name", "")
                row_cells[1].text = str(item.get("quantity", ""))
                row_cells[2].text = str(item.get("price", ""))
            
            output_path = f"F:\\AI系统开发原始框架\\output\\quotation.docx"
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            doc.save(output_path)
            
            return {
                "success": True,
                "message": f"已生成报价单: {output_path}",
                "result": {"path": output_path}
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_generate_contract(self, party_info: dict, terms: list) -> Dict[str, Any]:
        """生成合同"""
        try:
            from docx import Document
            doc = Document()
            doc.add_heading("装修合同", level=1)
            
            doc.add_heading("合同方信息", level=2)
            for key, value in party_info.items():
                doc.add_paragraph(f"{key}: {value}")
            
            doc.add_heading("合同条款", level=2)
            for i, term in enumerate(terms, 1):
                doc.add_paragraph(f"{i}. {term}")
            
            output_path = f"F:\\AI系统开发原始框架\\output\\contract.docx"
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            doc.save(output_path)
            
            return {
                "success": True,
                "message": f"已生成合同: {output_path}",
                "result": {"path": output_path}
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_recognize_material(self, image_path: str) -> Dict[str, Any]:
        """识别材料"""
        return {
            "success": False,
            "message": "材料识别功能需要视觉模型支持，暂未实现"
        }
    
    def _execute_recognize_handwriting(self, image_path: str) -> Dict[str, Any]:
        """识别手写"""
        return {
            "success": False,
            "message": "手写识别功能需要视觉模型支持，暂未实现"
        }
    
    def _execute_recognize_blueprint(self, image_path: str) -> Dict[str, Any]:
        """识别图纸"""
        return {
            "success": False,
            "message": "图纸识别功能需要视觉模型支持，暂未实现"
        }
    
    def _execute_online_design(self, room_info: dict, requirements: list) -> Dict[str, Any]:
        """在线设计"""
        return {
            "success": False,
            "message": "在线设计功能需要外部服务支持，暂未实现"
        }
    
    def _execute_generate_video(self, prompt: str, duration: int = 10) -> Dict[str, Any]:
        """生成视频"""
        return {
            "success": False,
            "message": "视频生成功能需要外部API支持，暂未实现"
        }
    
    def _execute_edit_video(self, video_path: str, operations: list) -> Dict[str, Any]:
        """编辑视频"""
        return {
            "success": False,
            "message": "视频编辑功能需要外部API支持，暂未实现"
        }
    
    def _execute_generate_social_content(self, platform: str, topic: str) -> Dict[str, Any]:
        """生成社交媒体内容"""
        return {
            "success": False,
            "message": "社交媒体内容生成功能需要外部API支持，暂未实现"
        }
    
    def _execute_web_search(self, query: str, search_engine: str = "baidu", max_results: int = 5) -> Dict[str, Any]:
        """联网搜索"""
        try:
            from services.web_search_tool import web_search_tool
            
            # 启动浏览器
            web_search_tool.start()
            
            # 执行搜索
            result = web_search_tool.search_web(query, search_engine, max_results)
            
            return {
                "success": True,
                "message": f"搜索完成: {query}",
                "result": result
            }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_chat(self, message: str) -> Dict[str, Any]:
        """对话"""
        try:
            from services.lightweight_model_service import lightweight_model_service
            response = lightweight_model_service.chat(
                model_name="qwen2-1.5b-int8",
                messages=[{"role": "user", "content": message}],
                temperature=0.7,
                max_tokens=512
            )
            
            if response.get("success"):
                return {
                    "success": True,
                    "message": response["response"],
                    "result": response
                }
            else:
                return {
                    "success": False,
                    "message": response.get("message", "对话失败")
                }
        except Exception as e:
            return {"success": False, "message": str(e)}
    
    def _execute_web_recognition(self, **kwargs) -> Dict[str, Any]:
        """智能网页识别"""
        try:
            from .web_recognition_tool import execute_web_recognition
            
            # 执行网页识别
            result = execute_web_recognition(**kwargs)
            
            return result
        except Exception as e:
            import traceback
            traceback.print_exc()
            return {"success": False, "message": str(e)}


# 创建全局实例
tool_executor = ToolExecutor()
